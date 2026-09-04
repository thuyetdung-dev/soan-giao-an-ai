import ast
import io
import json
import math
import operator
import re
import uuid
from dataclasses import dataclass
from typing import Any

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import streamlit as st
from docx import Document
import google.generativeai as genai
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from audit_engine import audit_exam, safe_autofix
from adaptive_engine import analyze_exam, variant_consistency, build_manifest
from pedagogy_engine import audit_pedagogy
from exam_factory import exam_generation_prompt, reviewer_prompt, parse_ai_json, certificate
from question_bank import QuestionBank, question_dna, fingerprint as question_fingerprint, select_from_bank
from v5_engine import build_variants, coverage_report, release_gate, manifest as build_v5_manifest
from lesson_engine import normalize_lesson, audit_lesson
from equation_engine import add_native_equation
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas

APP_VERSION = "7.0.0 (Math Premium + Exam Intelligence V5)"
MAX_UPLOAD_MB = 20
MAX_SOURCE_CHARS = 60_000
MAX_SLIDES = 60

THEMES = {
    "Xanh học thuật": {"primary": (16, 62, 105), "accent": (31, 150, 180), "light": (235, 246, 250)},
    "Xanh lá hiện đại": {"primary": (23, 92, 72), "accent": (45, 166, 116), "light": (235, 248, 242)},
    "Tím công nghệ": {"primary": (75, 53, 123), "accent": (137, 99, 186), "light": (245, 240, 251)},
    "Đỏ trang trọng": {"primary": (142, 32, 45), "accent": (210, 76, 89), "light": (253, 240, 242)},
}

ACTIVITY_COLORS = {
    "KHỞI ĐỘNG": (239, 108, 0),
    "HÌNH THÀNH KIẾN THỨC": (21, 101, 192),
    "LUYỆN TẬP": (0, 130, 100),
    "VẬN DỤNG": (126, 68, 153),
    "CỦNG CỐ": (180, 45, 55),
}

@dataclass
class LessonConfig:
    teacher: str
    school: str
    grade: str
    book: str
    lesson: str
    periods: int
    student_level: str
    slide_count: int
    theme_name: str
    include_answers: bool
    include_notes: bool

def clean_text(value: Any) -> str:
    text = str(value or "")
    text = re.sub(r"\s+", " ", text).strip()
    
    # 1. Khắc phục lỗi dư dấu gạch chéo ngược (R \\ {-2} -> R \ {-2})
    text = text.replace('\\\\', '\\')
    
    # 2. Định dạng Unicode toán học
    replacements = {"<=>": "⇔", "=>": "⇒", ">=": "≥", "<=": "≤", "+-": "±"}
    for old, new in replacements.items():
        text = text.replace(old, new)
        
    superscripts = str.maketrans("0123456789-+", "⁰¹²³⁴⁵⁶⁷⁸⁹⁻⁺")
    subscripts = str.maketrans("0123456789-+", "₀₁₂₃₄₅₆₇₈₉₋₊")
    text = re.sub(r"\^\{?(-?\d+)\}?", lambda m: m.group(1).translate(superscripts), text)
    text = re.sub(r"_\{?(-?\d+)\}?", lambda m: m.group(1).translate(subscripts), text)
    
    # 3. KEO TÀNG HÌNH BẢO VỆ CÔNG THỨC TOÁN
    # Dùng Non-Breaking Space (\xA0) thay cho khoảng trắng để PPT không bị cắt chữ nửa chừng
    math_ops = ['=', '≤', '≥', '∈', '∉', '≠', '⇔', '⇒', '±', '+', '-', '×', '÷', '<', '>']
    for op in math_ops:
        text = text.replace(f" {op} ", f"\xA0{op}\xA0")
    
    # Bảo vệ riêng ký hiệu hiệu tập hợp (VD: R \ {-2})
    text = text.replace(" \\ ", "\xA0\\\xA0")
    
    return text

def extract_docx(data: bytes) -> str:
    doc = Document(io.BytesIO(data))
    parts: list[str] = []
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            parts.append(paragraph.text.strip())
    for table_index, table in enumerate(doc.tables, 1):
        parts.append(f"[BẢNG {table_index}]")
        for row in table.rows:
            cells = [clean_text(cell.text) for cell in row.cells]
            parts.append(" | ".join(cells))
    return "\n".join(parts)

def read_source(uploaded_file) -> tuple[str, bytes, str]:
    data = uploaded_file.getvalue()
    if len(data) > MAX_UPLOAD_MB * 1024 * 1024:
        raise ValueError(f"Tệp vượt quá {MAX_UPLOAD_MB} MB.")
    name = uploaded_file.name.lower()
    if name.endswith(".docx"):
        return extract_docx(data)[:MAX_SOURCE_CHARS], data, "docx"
    if name.endswith(".txt"):
        try:
            text = data.decode("utf-8")
        except UnicodeDecodeError:
            text = data.decode("utf-8", errors="replace")
        return text[:MAX_SOURCE_CHARS], data, "txt"
    if name.endswith(".pdf"):
        return "", data, "pdf"
    if name.endswith(".json"):
        try:
            parsed = json.loads(data.decode("utf-8-sig"))
        except (UnicodeDecodeError, json.JSONDecodeError) as exc:
            raise ValueError("Tệp JSON không hợp lệ.") from exc
        return json.dumps(parsed, ensure_ascii=False)[:MAX_SOURCE_CHARS], data, "json"
    raise ValueError("Chỉ hỗ trợ PDF, DOCX, TXT và JSON.")

ALLOWED_FUNCTIONS = {
    "sin": np.sin, "cos": np.cos, "tan": np.tan, "sqrt": np.sqrt,
    "abs": np.abs, "exp": np.exp, "log": np.log, "ln": np.log,
}
ALLOWED_CONSTANTS = {"pi": np.pi, "e": np.e}
ALLOWED_BINOPS = (ast.Add, ast.Sub, ast.Mult, ast.Div, ast.Pow, ast.Mod)
ALLOWED_UNARYOPS = (ast.UAdd, ast.USub)
BINOP_FUNCS = {
    ast.Add: operator.add, ast.Sub: operator.sub, ast.Mult: operator.mul,
    ast.Div: operator.truediv, ast.Pow: operator.pow, ast.Mod: operator.mod,
}
UNARY_FUNCS = {ast.UAdd: operator.pos, ast.USub: operator.neg}

def _check_math_ast(node: ast.AST) -> None:
    if isinstance(node, ast.Expression):
        _check_math_ast(node.body)
    elif isinstance(node, ast.BinOp) and isinstance(node.op, ALLOWED_BINOPS):
        _check_math_ast(node.left)
        _check_math_ast(node.right)
    elif isinstance(node, ast.UnaryOp) and isinstance(node.op, ALLOWED_UNARYOPS):
        _check_math_ast(node.operand)
    elif isinstance(node, ast.Call):
        if not isinstance(node.func, ast.Name) or node.func.id not in ALLOWED_FUNCTIONS or node.keywords:
            raise ValueError("Hàm toán học không được hỗ trợ.")
        for arg in node.args:
            _check_math_ast(arg)
    elif isinstance(node, ast.Name):
        if node.id not in {"x", *ALLOWED_CONSTANTS}:
            raise ValueError(f"Biến '{node.id}' không được phép.")
    elif isinstance(node, ast.Constant):
        if not isinstance(node.value, (int, float)):
            raise ValueError("Hằng số không hợp lệ.")
    else:
        raise ValueError("Biểu thức chứa thành phần không an toàn.")

def _evaluate_math_ast(node: ast.AST, scope: dict[str, Any]) -> Any:
    if isinstance(node, ast.Expression):
        return _evaluate_math_ast(node.body, scope)
    if isinstance(node, ast.BinOp):
        left = _evaluate_math_ast(node.left, scope)
        right = _evaluate_math_ast(node.right, scope)
        return BINOP_FUNCS[type(node.op)](left, right)
    if isinstance(node, ast.UnaryOp):
        return UNARY_FUNCS[type(node.op)](_evaluate_math_ast(node.operand, scope))
    if isinstance(node, ast.Call):
        args = [_evaluate_math_ast(arg, scope) for arg in node.args]
        return scope[node.func.id](*args)
    if isinstance(node, ast.Name):
        return scope[node.id]
    if isinstance(node, ast.Constant):
        return node.value
    raise ValueError("Không thể tính biểu thức.")

def safe_math_eval(expression: str, x: np.ndarray) -> np.ndarray:
    normalized = expression.replace("^", "**").strip()
    tree = ast.parse(normalized, mode="eval")
    _check_math_ast(tree)
    scope = {"x": x, **ALLOWED_FUNCTIONS, **ALLOWED_CONSTANTS}
    result = _evaluate_math_ast(tree, scope)
    y = np.asarray(result, dtype=float)
    if y.ndim == 0:
        y = np.full_like(x, float(y))
    if y.shape != x.shape:
        raise ValueError("Biểu thức không tạo được một giá trị y cho mỗi x.")
    return y

def create_graph(graph: dict[str, Any]) -> io.BytesIO | None:
    expression = str(graph.get("expression", "x"))[:300]
    x_min = max(-100.0, min(100.0, float(graph.get("x_min", -5))))
    x_max = max(-100.0, min(100.0, float(graph.get("x_max", 5))))
    if not x_min < x_max:
        return None
    x = np.linspace(x_min, x_max, 1600)
    with np.errstate(all="ignore"):
        y = safe_math_eval(expression, x)
    finite = np.isfinite(y)
    if not finite.any():
        return None
    finite_values = y[finite]
    q1, q99 = np.percentile(finite_values, [1, 99])
    span = max(4.0, q99 - q1)
    lower, upper = q1 - span * 0.12, q99 + span * 0.12
    plot_y = y.copy()
    jumps = np.zeros_like(plot_y, dtype=bool)
    jumps[1:] = np.abs(np.diff(plot_y)) > span * 0.8
    plot_y[~finite | jumps | (plot_y < lower - span) | (plot_y > upper + span)] = np.nan

    fig, ax = plt.subplots(figsize=(8.4, 4.6))
    ax.plot(x, plot_y, color="#1261a0", linewidth=2.4)
    ax.axhline(0, color="#263238", linewidth=1)
    ax.axvline(0, color="#263238", linewidth=1)
    ax.grid(True, linestyle="--", alpha=0.28)
    ax.set_xlim(x_min, x_max)
    ax.set_ylim(lower, upper)
    ax.set_xlabel("x")
    ax.set_ylabel("y", rotation=0, labelpad=10)
    ax.set_title(clean_text(graph.get("caption") or f"Đồ thị y = {expression}"), fontsize=13, weight="bold")
    fig.tight_layout()
    output = io.BytesIO()
    fig.savefig(output, format="png", dpi=220, transparent=False, facecolor="white")
    plt.close(fig)
    output.seek(0)
    return output

def create_variation_table(data: dict[str, Any]) -> io.BytesIO | None:
    points = data.get("points", [])
    intervals = data.get("interval_signs", [])
    values = data.get("values", [])
    if not isinstance(points, list) or len(points) < 2 or len(points) > 12:
        return None
    cols = 2 * len(points) - 1
    fig, ax = plt.subplots(figsize=(max(7.5, cols * 0.72), 3.2))
    ax.set_xlim(0, cols + 1)
    ax.set_ylim(0, 3)
    ax.axis("off")
    for y in range(4):
        ax.plot([0, cols + 1], [y, y], color="#263238", lw=1.1)
    ax.plot([1, 1], [0, 3], color="#263238", lw=1.2)
    ax.text(.5, 2.5, "x", ha="center", va="center", fontsize=14, style="italic")
    ax.text(.5, 1.5, "y′", ha="center", va="center", fontsize=14, style="italic")
    ax.text(.5, .5, "y", ha="center", va="center", fontsize=14, style="italic")
    for i, point in enumerate(points):
        xpos = 1.5 + 2 * i
        ax.text(xpos, 2.5, clean_text(point), ha="center", va="center", fontsize=13)
        if i < len(values) and values[i] not in (None, ""):
            val = clean_text(values[i])
            ypos = .78 if i == 0 else (.22 if i == len(points) - 1 else .5)
            ax.text(xpos, ypos, val, ha="center", va="center", fontsize=12)
        if i < len(points) - 1:
            sign = clean_text(intervals[i]) if i < len(intervals) else ""
            mid = xpos + 1
            ax.text(mid, 1.5, sign, ha="center", va="center", fontsize=14)
            rising = sign == "+"
            ax.annotate("", xy=(xpos + 1.72, .78 if rising else .22),
                        xytext=(xpos + .28, .22 if rising else .78),
                        arrowprops={"arrowstyle": "->", "lw": 1.6, "color": "#1261a0"})
    fig.tight_layout(pad=.3)
    output = io.BytesIO()
    fig.savefig(output, format="png", dpi=220, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    output.seek(0)
    return output

def split_bullets(items: list[str], max_chars: int = 430) -> list[list[str]]:
    pages: list[list[str]] = []
    current: list[str] = []
    size = 0
    for raw in items:
        item = clean_text(raw)
        if not item:
            continue
        if current and (size + len(item) > max_chars or len(current) >= 6):
            pages.append(current)
            current, size = [], 0
        current.append(item)
        size += len(item)
    if current or not pages:
        pages.append(current)
    return pages

def validate_lesson(data: Any) -> dict[str, Any]:
    return normalize_lesson(data, MAX_SLIDES)

def build_prompt(config: LessonConfig) -> str:
    answer_rule = "Có thể cung cấp đáp án ở trường answer." if config.include_answers else "Trường answer luôn để trống."
    lesson_name = config.lesson.strip() or "Tự xác định chính xác từ tài liệu nguồn"
    return f"""
Bạn là chuyên gia thiết kế bài giảng Toán theo Chương trình GDPT 2018.
THÔNG TIN BẮT BUỘC:
- Khối: {config.grade}; bộ sách: {config.book}; bài: {lesson_name}.
- Thời lượng: {config.periods} tiết; đối tượng: {config.student_level}.
- Tạo khoảng {config.slide_count} slide, phân bổ hợp lý theo đúng thời lượng.

NGUYÊN TẮC NỘI DUNG:
- Xác định đúng chủ đề từ tên bài và tài liệu nguồn; tuyệt đối không đưa nội dung của một bài khác.
- Tài liệu nguồn là căn cứ chính. Không bịa định nghĩa, định lý, dữ kiện, ví dụ hay số trang.
- Nếu nguồn thiếu thông tin thiết yếu, nêu rõ trong teacher_note thay vì tự khẳng định.
- Tổ chức theo các hoạt động: KHỞI ĐỘNG, HÌNH THÀNH KIẾN THỨC, LUYỆN TẬP, VẬN DỤNG, CỦNG CỐ.
- Mỗi hoạt động phải có nhiệm vụ rõ, sản phẩm mong đợi và cách kiểm tra nhanh phù hợp học sinh {config.student_level}.
- Mỗi định nghĩa, quy tắc, ví dụ hoặc câu hỏi chính đặt trên slide độc lập; tối đa 5–6 dòng/slide. {answer_rule}
- Công thức trong formulas viết bằng LaTeX chuẩn, không đặt dấu $ bao quanh. Hỗ trợ tốt \\frac, \\sqrt, số mũ, chỉ số, chữ Hy Lạp, tổng, tích phân và các quan hệ. Không dùng ký hiệu mơ hồ hoặc kết luận thiếu điều kiện.
- Dùng ít nhất 6 kiểu layout phù hợp nội dung, không lặp một kiểu quá 4 slide liên tiếp.
- Chuỗi slide phải theo tiến trình học tập: vấn đề → khám phá → khái quát → luyện tập → vận dụng → tự đánh giá.

ĐỒ HỌA TÙY CHỌN:
- graph: {{"expression":"x**3-3*x", "x_min":-5, "x_max":5, "caption":"..."}}. Chỉ dùng Python math chuẩn (sin, cos, exp).
- variation_table: {{"points":["-∞","-1","1","+∞"], "interval_signs":["+","-","+"], "values":["-∞","2","-2","+∞"]}}.

Trả về duy nhất JSON chuẩn:
{{
  "title":"Tên bài học",
  "objectives":["Yêu cầu cần đạt 1","Yêu cầu cần đạt 2"],
  "slides":[{{
    "title":"Tiêu đề slide",
    "subtitle":"Thông điệp ngắn nếu cần",
    "activity":"HÌNH THÀNH KIẾN THỨC",
    "layout":"section|concept|process|example|practice|compare|visual|quiz|summary|content",
    "bullets":["Ý 1","Ý 2"],
    "formulas":["f'(x)=0"],
    "question":"Câu hỏi hoặc nhiệm vụ dành cho học sinh",
    "product":"Sản phẩm học tập mong đợi",
    "answer":"",
    "teacher_note":"Ghi chú giảng dạy chi tiết cho giáo viên",
    "source_ref":"Tên tài liệu, mục hoặc trang làm căn cứ cho slide",
    "graph":null,
    "variation_table":null
  }}]
}}
"""

def generate_lesson(model_name: str, source_text: str, source_bytes: bytes, source_type: str, config: LessonConfig) -> dict[str, Any]:
    prompt = build_prompt(config)
    if source_type == "pdf":
        contents = [{"mime_type": "application/pdf", "data": source_bytes}, prompt]
    else:
        contents = [f"TÀI LIỆU NGUỒN:\n{source_text}\n\n{prompt}"]
    
    model = genai.GenerativeModel(model_name, generation_config={"response_mime_type": "application/json", "temperature": 0.25})
    response = model.generate_content(contents)
    
    if not response.text:
        raise ValueError("AI không trả về nội dung. Vui lòng thử lại.")
        
    raw_json = response.text
    raw_json = re.sub(r'```(?:json)?', '', raw_json).strip()
    raw_json = re.sub(r'```', '', raw_json).strip()
    
    try:
        match = re.search(r'\{.*\}', raw_json, re.DOTALL)
        if match:
            return validate_lesson(json.loads(match.group(0), strict=False))
        return validate_lesson(json.loads(raw_json))
    except Exception:
        raise ValueError("AI trả về JSON không hợp lệ. Vui lòng tạo lại; hệ thống không tự biến đổi công thức để tránh làm sai nội dung.")

def add_full_background(slide, color: tuple[int, int, int]) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)

def add_text(slide, text: str, left: float, top: float, width: float, height: float,
             size: int, color=(40, 40, 40), bold=False, align=PP_ALIGN.LEFT,
             font="Aptos", valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(.08)
    frame.margin_top = frame.margin_bottom = Inches(.04)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)
    return box

def add_header(slide, title: str, activity: str, theme: dict[str, tuple[int, int, int]], page: int):
    primary = theme["primary"]
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(.18))
    bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(*primary); bar.line.fill.background()
    add_text(slide, title, .72, .38, 11.8, .7, 27, primary, True)
    activity_color = ACTIVITY_COLORS.get(activity, theme["accent"])
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(.75), Inches(1.08), Inches(3.1), Inches(.42))
    pill.fill.solid(); pill.fill.fore_color.rgb = RGBColor(*activity_color); pill.line.fill.background()
    add_text(slide, activity, .82, 1.13, 2.95, .28, 11, (255, 255, 255), True, PP_ALIGN.CENTER)
    add_text(slide, str(page), 12.35, 7.05, .4, .25, 10, (100, 100, 100), False, PP_ALIGN.RIGHT)

def add_bullets(slide, bullets: list[str], left: float, top: float, width: float, height: float,
                font_size: int = 21):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    # Kích hoạt Word Wrap của PPT
    frame.word_wrap = True
    frame.margin_left = Inches(.08); frame.margin_right = Inches(.08)
    
    # Gộp toàn bộ đoạn văn vào CHUNG MỘT Textbox để hỗ trợ hiệu ứng (Animation By Paragraph)
    for i, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = clean_text(bullet)
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(43, 48, 53)
        p.space_after = Pt(10)
        p.text = "●  " + p.text
    return box

def add_answer_box(slide, answer: str, theme):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(.82), Inches(5.72), Inches(11.7), Inches(.92))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(*theme["light"])
    shape.line.color.rgb = RGBColor(*theme["accent"])
    add_text(slide, "ĐÁP ÁN/GỢI Ý: " + answer, 1.0, 5.9, 11.3, .55, 16, theme["primary"], True)


def add_panel(slide, left, top, width, height, fill, line=None, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(*fill)
    if line: shape.line.color.rgb = RGBColor(*line)
    else: shape.line.fill.background()
    return shape


def add_formula_block(slide, formulas: list[str], left: float, top: float, width: float, theme):
    if not formulas: return
    height = min(1.35, .42 + .38 * len(formulas))
    add_panel(slide, left, top, width, height, theme["light"], theme["accent"])
    row_height=(height-.16)/max(1,len(formulas))
    for index,formula in enumerate(formulas):
        add_native_equation(slide,formula,left+.16,top+.08+index*row_height,width-.32,row_height,22,theme["primary"])


def add_learning_task(slide, slide_data, theme, top=4.85):
    question=slide_data.get("question",""); product=slide_data.get("product","")
    if question:
        add_panel(slide,.82,top,7.6,.72,theme["light"],theme["accent"])
        add_text(slide,"NHIỆM VỤ  •  "+question,1.02,top+.13,7.2,.43,16,theme["primary"],True)
    if product:
        add_panel(slide,8.65,top,3.87,.72,(248,249,250),(205,211,217))
        add_text(slide,"SẢN PHẨM  •  "+product,8.83,top+.13,3.5,.43,14,(55,65,72),True)


def add_native_variation_table(slide, data, left, top, width, height, theme):
    points=data.get("points",[]); signs=data.get("interval_signs",[]); values=data.get("values",[])
    if not isinstance(points,list) or len(points)<2: return False
    cols=2*len(points)
    table_shape=slide.shapes.add_table(3,cols,Inches(left),Inches(top),Inches(width),Inches(height))
    table=table_shape.table
    first_width=.8; table.columns[0].width=Inches(first_width)
    remaining=max(.3,(width-first_width)/(cols-1))
    for col in range(1,cols): table.columns[col].width=Inches(remaining)
    labels=["x","y′","y"]
    for row in range(3):
        for col in range(cols):
            cell=table.cell(row,col); cell.fill.solid(); cell.fill.fore_color.rgb=RGBColor(255,255,255)
            cell.margin_left=cell.margin_right=Inches(.03); cell.margin_top=cell.margin_bottom=Inches(.02)
            p=cell.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
            run=p.add_run(); run.font.name="Cambria Math"; run.font.size=Pt(17); run.font.color.rgb=RGBColor(35,45,52)
            if col==0: run.text=labels[row]; run.font.bold=True; cell.fill.fore_color.rgb=RGBColor(*theme["light"])
            elif col%2==1:
                idx=(col-1)//2
                if row==0: run.text=str(points[idx]) if idx<len(points) else ""
                elif row==2: run.text=str(values[idx]) if idx<len(values) else ""
            else:
                idx=col//2-1; sign=str(signs[idx]) if idx<len(signs) else ""
                if row==1: run.text=sign
                elif row==2: run.text="↗" if sign=="+" else "↘" if sign=="-" else "→"
    return True

def build_pptx(lesson: dict[str, Any], config: LessonConfig) -> bytes:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]
    theme = THEMES[config.theme_name]

    cover = prs.slides.add_slide(blank)
    add_full_background(cover, theme["primary"])
    accent = cover.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(.85), Inches(1.0), Inches(.12), Inches(5.4))
    accent.fill.solid(); accent.fill.fore_color.rgb = RGBColor(*theme["accent"]); accent.line.fill.background()
    add_text(cover, config.grade.upper() + " • " + config.book.upper(), 1.28, 1.1, 10.8, .42, 15, theme["light"], True)
    add_text(cover, lesson.get("title") or config.lesson, 1.25, 1.75, 10.8, 2.25, 38, (255, 255, 255), True, valign=MSO_ANCHOR.MIDDLE)
    info = f"{config.periods} tiết  •  Giáo viên: {config.teacher or '................................'}"
    if config.school:
        info += f"\n{config.school}"
    add_text(cover, info, 1.28, 4.55, 10.5, 1.0, 18, theme["light"])
    add_text(cover, "POWERPOINT BÀI GIẢNG TOÁN THPT", 1.28, 6.65, 10.5, .35, 12, theme["light"], True)

    if lesson.get("objectives"):
        objective_slide=prs.slides.add_slide(blank); add_full_background(objective_slide,(255,255,255))
        add_header(objective_slide,"Sau bài học, học sinh làm được gì?","KHỞI ĐỘNG",theme,1)
        add_bullets(objective_slide,lesson["objectives"],1.05,1.85,11.2,4.7,22)

    page = 2 if lesson.get("objectives") else 1
    for slide_data in lesson["slides"]:
        chunks = split_bullets(slide_data["bullets"])
        for chunk_index, chunk in enumerate(chunks):
            slide = prs.slides.add_slide(blank)
            layout=slide_data.get("layout","content")
            add_full_background(slide, theme["primary"] if layout=="section" else (255, 255, 255))
            title = slide_data["title"] + (f" ({chunk_index + 1}/{len(chunks)})" if len(chunks) > 1 else "")
            activity = slide_data["activity"]
            if layout=="section":
                add_text(slide,activity,.95,1.0,11.4,.45,15,theme["light"],True)
                add_text(slide,title,.95,1.72,11.2,2.0,38,(255,255,255),True,valign=MSO_ANCHOR.MIDDLE)
                if slide_data.get("subtitle"): add_text(slide,slide_data["subtitle"],.98,4.05,10.9,.85,20,theme["light"])
                add_text(slide,str(page),12.2,6.85,.45,.3,11,theme["light"],False,PP_ALIGN.RIGHT)
                if config.include_notes and (slide_data.get("teacher_note") or slide_data.get("source_ref")):
                    slide.notes_slide.notes_text_frame.text=(slide_data.get("teacher_note","")+"\n\n[Sources]\n- "+(slide_data.get("source_ref") or "Tài liệu nguồn người dùng cung cấp")).strip()
                page += 1
                continue
            add_header(slide, title, activity, theme, page)
            visual = None
            if chunk_index == 0 and slide_data.get("graph"):
                try:
                    visual = create_graph(slide_data["graph"])
                except (ValueError, SyntaxError, TypeError):
                    visual = None
            native_variation = chunk_index == 0 and bool(slide_data.get("variation_table"))

            has_answer = config.include_answers and bool(slide_data.get("answer"))
            content_bottom = 5.48 if has_answer else 6.78
            formulas=slide_data.get("formulas",[])

            if native_variation:
                add_bullets(slide,chunk,.82,1.65,11.7,1.2,18)
                add_native_variation_table(slide,slide_data["variation_table"],.82,2.85,11.7,2.45,theme)
            elif layout=="process" and len(chunk)>=2:
                step_w=11.55/min(4,len(chunk))
                for idx,item in enumerate(chunk[:4]):
                    left=.82+idx*step_w
                    add_panel(slide,left,1.85,step_w-.22,2.5,theme["light"],theme["accent"])
                    circle=slide.shapes.add_shape(MSO_SHAPE.OVAL,Inches(left+.12),Inches(1.96),Inches(.68),Inches(.68)); circle.fill.solid(); circle.fill.fore_color.rgb=RGBColor(*theme["accent"]); circle.line.fill.background()
                    add_text(slide,str(idx+1),left+.18,2.03,.55,.55,25,(255,255,255),True,PP_ALIGN.CENTER,valign=MSO_ANCHOR.MIDDLE)
                    add_text(slide,item,left+.18,2.82,step_w-.55,1.25,17,(43,48,53),True,PP_ALIGN.CENTER)
                add_formula_block(slide,formulas,2.25,4.62,8.8,theme)
            elif layout in {"concept","example","compare"} and not visual:
                add_panel(slide,.82,1.75,7.25,3.65,(255,255,255),(215,221,226))
                add_bullets(slide,chunk,1.04,2.02,6.8,3.05,20)
                add_panel(slide,8.35,1.75,4.17,3.65,theme["light"],theme["accent"])
                side_text=slide_data.get("question") or slide_data.get("subtitle") or slide_data.get("product") or "Quan sát, trao đổi và rút ra kết luận."
                add_text(slide,side_text,8.68,2.08,3.5,2.25,20,theme["primary"],True,PP_ALIGN.CENTER,valign=MSO_ANCHOR.MIDDLE)
                add_formula_block(slide,formulas,8.62,4.35,3.63,theme)
            elif visual:
                # ÉP CHIỀU RỘNG CHỮ TỐI ĐA LÀ 7.0 INCH (~1/2 trang theo chiều ngang)
                # Chữ chạy tới mốc này sẽ TỰ ĐỘNG rớt dòng, hoàn toàn không che lấp Hình bên phải.
                add_bullets(slide, chunk, .75, 1.72, 7.0, content_bottom - 1.72, 21)
                
                # Cố định hình vẽ/bảng biến thiên ở góc trên bên phải
                slide.shapes.add_picture(visual, Inches(8.2), Inches(1.08), width=Inches(4.8))
            else:
                size = 23 if sum(map(len, chunk)) < 280 else 21
                body_bottom=4.72 if slide_data.get("question") or slide_data.get("product") else content_bottom
                add_bullets(slide, chunk, .82, 1.72, 11.7, body_bottom - 1.72, size)
                if formulas: add_formula_block(slide,formulas,2.15,max(3.75,body_bottom-.9),9.0,theme)

            if (slide_data.get("question") or slide_data.get("product")) and layout not in {"concept","example","compare"}:
                add_learning_task(slide,slide_data,theme,4.78 if not has_answer else 4.72)
                
            if has_answer:
                add_answer_box(slide, slide_data["answer"], theme)
            if config.include_notes and (slide_data.get("teacher_note") or slide_data.get("source_ref")):
                notes_frame = slide.notes_slide.notes_text_frame
                notes_frame.text = (slide_data.get("teacher_note","")+"\n\n[Sources]\n- "+(slide_data.get("source_ref") or "Tài liệu nguồn người dùng cung cấp")).strip()
            page += 1

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


def build_exam_docx(exam: dict[str, Any], certificate_data: dict[str, Any] | None = None) -> bytes:
    doc = Document()
    doc.add_heading(exam.get("title", "ĐỀ KIỂM TRA TOÁN"), 0)
    doc.add_paragraph(f"{exam.get('subject','Toán')} • {exam.get('grade','')}" )
    for i, q in enumerate(exam.get("questions", []), 1):
        doc.add_heading(f"Câu {i}. {q.get('question','')}", level=2)
        typ = str(q.get("type", "mcq"))
        if typ in {"mcq", "multiple_choice"}:
            for j, opt in enumerate(q.get("options", [])):
                doc.add_paragraph(f"{chr(65+j)}. {opt}")
        elif typ in {"true_false", "tf"}:
            for j, stt in enumerate(q.get("statements", q.get("options", []))):
                doc.add_paragraph(f"{j+1}) {stt}")
        else:
            doc.add_paragraph("Trả lời: ........................................................")
    if certificate_data:
        doc.add_page_break()
        doc.add_heading("BÁO CÁO QA V5.0", 1)
        doc.add_paragraph(json.dumps(certificate_data, ensure_ascii=False, indent=2))
    out=io.BytesIO(); doc.save(out); return out.getvalue()


def build_exam_pptx(exam: dict[str, Any], theme_name: str = "Xanh học thuật", show_answers: bool = False) -> bytes:
    prs=Presentation(); prs.slide_width, prs.slide_height=Inches(13.333), Inches(7.5); blank=prs.slide_layouts[6]
    theme=THEMES[theme_name]
    cover=prs.slides.add_slide(blank); add_full_background(cover, theme["primary"])
    add_text(cover, "AI EXAM FACTORY V5.0", 1.0, 1.0, 11.0, .5, 18, theme["light"], True)
    add_text(cover, exam.get("title", "ĐỀ KIỂM TRA TOÁN"), 1.0, 2.0, 11.2, 2.0, 34, (255,255,255), True, valign=MSO_ANCHOR.MIDDLE)
    add_text(cover, f"{exam.get('grade','')} • {exam.get('subject','Toán')}", 1.0, 4.5, 10.5, .5, 18, theme["light"])
    for i,q in enumerate(exam.get("questions", []),1):
        slide=prs.slides.add_slide(blank); add_full_background(slide,(255,255,255))
        add_header(slide, f"Câu {i}", str(q.get("level", "")).upper() or "ĐỀ KIỂM TRA", theme, i)
        bullets=[clean_text(q.get("question",""))]
        typ=q.get("type","mcq")
        if typ in {"mcq","multiple_choice"}:
            bullets += [f"{chr(65+j)}. {clean_text(o)}" for j,o in enumerate(q.get("options",[]))]
        elif typ in {"true_false","tf"}:
            bullets += [f"{j+1}) {clean_text(o)}" for j,o in enumerate(q.get("statements",q.get("options",[])))]
        add_bullets(slide, bullets, .8, 1.7, 11.7, 4.6, 19)
        if show_answers and q.get("solution"):
            add_answer_box(slide, clean_text(q.get("solution")), theme)
    out=io.BytesIO(); prs.save(out); return out.getvalue()


def combined_v4_report(exam: dict[str, Any], math_report: dict, ped_report: dict, council: list[dict]) -> dict[str, Any]:
    council_fail=any(str(x.get("status")).upper()=="FAIL" for x in council)
    council_review=any(str(x.get("status")).upper() in {"REVIEW","MANUAL_REVIEW"} for x in council)
    status="FAIL" if math_report["status"]=="FAIL" or ped_report["status"]=="FAIL" or council_fail else "MANUAL_REVIEW" if math_report["status"]=="MANUAL_REVIEW" or ped_report["status"]=="MANUAL_REVIEW" or council_review else "PASS"
    scores=[math_report["summary"]["score"], ped_report["summary"]["score"]]+[float(x.get("score",0)) for x in council if isinstance(x.get("score"),(int,float))]
    score=round(sum(scores)/len(scores),1) if scores else 0
    base={"version":"5.0.0","status":status,"summary":{"total":len(exam.get("questions",[])),"score":score},"math":math_report,"pedagogy":ped_report,"council":council}
    base["certificate"]=certificate(base)
    return base


def validate_blueprint_counts(total: int, type_counts: dict[str, int], level_counts: dict[str, int]) -> None:
    if sum(type_counts.values()) != total:
        raise ValueError(f"Tổng các loại câu phải bằng {total}; hiện là {sum(type_counts.values())}.")
    if sum(level_counts.values()) != total:
        raise ValueError(f"Tổng phân bố mức độ phải bằng {total}; hiện là {sum(level_counts.values())}.")


def balanced_topic_distribution(raw_topics: str, total: int) -> dict[str, int]:
    topics = [line.strip() for line in raw_topics.splitlines() if line.strip()]
    if not topics:
        raise ValueError("Cần nhập ít nhất một chủ đề.")
    base, remainder = divmod(total, len(topics))
    return {topic: base + (1 if index < remainder else 0) for index, topic in enumerate(topics)}

def get_api_key() -> str:
    try:
        return st.secrets.get("GEMINI_API_KEY", "")
    except Exception:
        return ""

st.set_page_config(page_title="Trợ lý PowerPoint Toán THPT", page_icon="📐", layout="wide")
st.markdown("""
<style>
.block-container{padding-top:1.4rem;max-width:1250px}.stButton>button{font-weight:700;border-radius:10px}
[data-testid="stSidebar"]{background:#f5f8fb}.small-note{color:#59636e;font-size:.92rem}
</style>
""", unsafe_allow_html=True)
st.title("📐 Trợ lý soạn PowerPoint Toán THPT")
st.caption(f"Phiên bản {APP_VERSION} • Đồ họa Toán học AST • Chống rớt chữ")

api_key = get_api_key()
with st.sidebar:
    mode = st.radio("Chế độ làm việc", ["Tạo bài giảng PowerPoint", "🏭 AI Exam Factory V5.0", "🧬 Exam Intelligence V5.0", "Thẩm định đề Toán Pro", "Thẩm định đề Toán 360°"], index=0)
if not api_key and mode not in {"Thẩm định đề Toán Pro", "Thẩm định đề Toán 360°", "🧬 Exam Intelligence V5.0"}:
    st.error("Chưa cấu hình GEMINI_API_KEY trong Secrets. Chế độ Thẩm định đề Toán Pro vẫn chạy offline không cần API.")
    st.stop()

with st.sidebar:
    st.header("Thông tin bài dạy")
    teacher = st.text_input("Tên giáo viên", placeholder="Nguyễn Văn A")
    school = st.text_input("Trường", placeholder="Trường THPT …")
    grade = st.selectbox("Khối lớp", ["Toán 10", "Toán 11", "Toán 12"])
    book = st.selectbox("Bộ sách", ["Kết nối tri thức", "Cánh Diều", "Chân trời sáng tạo", "Tài liệu riêng"])
    lesson = st.text_input("Tên bài", placeholder="Có thể để trống để AI xác định")
    periods = st.number_input("Số tiết", 1, 6, 2)
    student_level = st.selectbox("Đối tượng học sinh", ["Trung bình – khá", "Đồng đều", "Còn hạn chế", "Khá – giỏi"])
    slide_count = st.slider("Số slide dự kiến", 15, 50, 30)
    theme_name = st.selectbox("Phong cách", list(THEMES))
    include_answers = st.checkbox("Đưa đáp án/gợi ý vào slide", True)
    include_notes = st.checkbox("Tạo ghi chú dành cho giáo viên", True)
    
    st.markdown("---")
    try:
        genai.configure(api_key=api_key)
        available_models = [m.name for m in genai.list_models() if 'generateContent' in m.supported_generation_methods]
        if available_models:
            default_index = available_models.index('models/gemini-1.5-flash') if 'models/gemini-1.5-flash' in available_models else 0
            selected_model = st.selectbox("🤖 Chọn mô hình AI:", available_models, index=default_index)
        else:
            st.error("Khóa API không có quyền truy cập.")
            selected_model = None
    except Exception:
        selected_model = "models/gemini-1.5-flash"

st.subheader("1. Tải tài liệu nguồn")
uploaded = st.file_uploader("PDF, Word, TXT hoặc JSON (tối đa 20 MB)", type=["pdf", "docx", "txt", "json"])
st.markdown('<div class="small-note">Nên dùng tài liệu chính thống: SGK, SGV, kế hoạch bài dạy hoặc chuyên đề đã kiểm duyệt.</div>', unsafe_allow_html=True)

if mode == "🧬 Exam Intelligence V5.0":
    st.subheader("🧬 Exam Intelligence V5.0 — Question DNA + Ngân hàng câu hỏi + Adaptive Exam Factory")
    st.info("V5.0: Ngân hàng câu hỏi → Question DNA → chọn theo ma trận → Math Engine + Sư phạm → Hội đồng 3 AI → nhiều mã đề → QA liên mã → Release Gate.")
    bank=QuestionBank("question_bank_v5.sqlite3")
    tabA,tabB,tabC=st.tabs(["🗃️ Ngân hàng câu hỏi","🏭 Tạo đề từ ngân hàng","📊 Analytics & DNA"])
    with tabA:
        st.write(f"**Số câu đang lưu:** {bank.count()}")
        bank_file=st.file_uploader("Nhập đề JSON để đưa vào ngân hàng",type=["json"],key="v5_bank_upload")
        if st.button("➕ NẠP CÂU HỎI VÀO NGÂN HÀNG",use_container_width=True) and bank_file:
            try:
                ex=json.loads(bank_file.getvalue().decode("utf-8")); st.success(bank.add_exam(ex))
            except Exception as e: st.error(str(e))
        c1,c2,c3=st.columns(3)
        topic_filter=c1.text_input("Lọc chủ đề",key="v5_topic")
        level_filter=c2.text_input("Lọc mức độ",key="v5_level")
        type_filter=c3.selectbox("Lọc loại",["","mcq","true_false","short_answer"],key="v5_type")
        rows=bank.search(topic_filter,level_filter,type_filter,100)
        st.dataframe([{"ID":q.get("id"),"Loại":qtype(q),"Mức độ":q.get("level"),"Chủ đề":q.get("topic"),"DNA":q.get("_dna",{}).get("fingerprint","")[:12],"Lần dùng":q.get("_uses",0)} for q in rows],use_container_width=True)
    with tabB:
        total_v5=st.number_input("Tổng số câu",1,60,10,key="v5_total")
        mcq_v5=st.number_input("MCQ",0,int(total_v5),min(10,int(total_v5)),key="v5_mcq")
        tf_v5=st.number_input("Đúng/Sai",0,int(total_v5),0,key="v5_tf")
        short_v5=st.number_input("Trả lời ngắn",0,int(total_v5),0,key="v5_short")
        levels_v5=st.text_input("Mức độ", "nhận biết:3, thông hiểu:3, vận dụng:3, vận dụng cao:1",key="v5_levels")
        topics_v5=st.text_area("Chủ đề ưu tiên", "Tính đơn điệu và cực trị của hàm số",key="v5_topics")
        codes_v5=st.number_input("Số mã đề",1,20,4,key="v5_codes")
        if st.button("🚀 TẠO ĐỀ V5.0 + RELEASE GATE",type="primary",use_container_width=True):
            try:
                bp={"total_questions":int(total_v5),"type_distribution":{"mcq":int(mcq_v5),"true_false":int(tf_v5),"short_answer":int(short_v5)},"level_distribution":{}}
                for item in levels_v5.split(','):
                    if ':' in item:
                        k,v=item.split(':',1); bp["level_distribution"][k.strip()]=int(v.strip())
                validate_blueprint_counts(int(total_v5), bp["type_distribution"], bp["level_distribution"])
                bp["topic_distribution"]=balanced_topic_distribution(topics_v5, int(total_v5))
                exam={"title":"Đề Toán V5.0","subject":"Toán","grade":grade,"blueprint":bp,"questions":select_from_bank(bank,bp)}
                if len(exam["questions"])<int(total_v5): st.warning(f"Ngân hàng chỉ đáp ứng {len(exam['questions'])}/{int(total_v5)} câu theo bộ lọc. Có thể bổ sung câu hoặc dùng AI Factory V5.0 để sinh câu mới.")
                mr=audit_exam(exam); pr=audit_pedagogy(exam)
                variants=build_variants(exam,int(codes_v5)); vc=variant_consistency(variants)
                gate=release_gate(mr,pr,[],vc); mf=build_v5_manifest(exam,variants,gate)
                bank.mark_used(exam["questions"])
                st.session_state["v5_exam"]=exam; st.session_state["v5_mr"]=mr; st.session_state["v5_pr"]=pr; st.session_state["v5_variants"]=variants; st.session_state["v5_manifest"]=mf
            except Exception as e: st.error(f"V5.0 lỗi: {e}")
        exam=st.session_state.get("v5_exam"); mr=st.session_state.get("v5_mr"); pr=st.session_state.get("v5_pr"); variants=st.session_state.get("v5_variants",[]); mf=st.session_state.get("v5_manifest",{})
        if exam:
            gate=mf.get("release_gate","CONDITIONAL"); a,b,c,d=st.columns(4); a.metric("Câu",len(exam.get("questions",[]))); b.metric("Gate",gate); c.metric("Mã đề",len(variants)); d.metric("DNA unique",coverage_report(exam).get("unique",0))
            if gate=="CERTIFIED": st.success("🟢 CERTIFIED — đề vượt Release Gate V5.0 trong phạm vi các bộ máy tự động.")
            elif gate=="REJECTED": st.error("🔴 REJECTED — cần sửa lỗi trước khi phát hành.")
            else: st.warning("🟡 CONDITIONAL — cần giáo viên duyệt các điểm REVIEW.")
            st.download_button("📥 Tải manifest V5.0",json.dumps(mf,ensure_ascii=False,indent=2),"manifest_v5.json","application/json",use_container_width=True)
            st.download_button("📥 Tải đề gốc JSON",json.dumps(exam,ensure_ascii=False,indent=2),"de_goc_v5.json","application/json",use_container_width=True)
    with tabC:
        st.json(bank.stats())
        qshow=bank.search(limit=20)
        if qshow:
            with st.expander("🧬 DNA của 20 câu gần nhất"):
                st.json([question_dna(q) for q in qshow])
    bank.close()
    st.stop()

if mode == "🏭 AI Exam Factory V5.0":
    st.subheader("🏭 AI Exam Factory V5.0 — AI Exam Factory + Adaptive Intelligence")
    st.info("Luồng V5.0: Ma trận → AI sinh đề → Math Engine → Sư phạm → 3 AI phản biện → phân tích độ khó → nhiều mã đề → QA liên mã → báo cáo kiểm định.")
    col1,col2=st.columns(2)
    with col1:
        total_q=st.number_input("Tổng số câu", 1, 60, 10)
        grade_exam=st.selectbox("Khối", ["Toán 10","Toán 11","Toán 12"], key="exam_grade")
        topics=st.text_area("Chủ đề (mỗi dòng một chủ đề)", "Tính đơn điệu và cực trị của hàm số\nGTLN – GTNN\nHàm số mũ và logarit")
    with col2:
        mcq_n=st.number_input("Số MCQ",0,int(total_q),max(1,min(10,int(total_q))))
        tf_n=st.number_input("Số Đúng/Sai",0,int(total_q),0)
        short_n=st.number_input("Số trả lời ngắn",0,int(total_q),0)
        levels=st.text_input("Phân bố mức độ", "nhận biết:3, thông hiểu:3, vận dụng:3, vận dụng cao:1")
        variant_n=st.number_input("Số mã đề", 1, 8, 4, key="variant_n")
        variant_strategy=st.selectbox("Chiến lược mã đề", ["Đảo thứ tự câu + phương án", "Biến thể tham số (AI sinh lại + kiểm chứng)"])
    if st.button("🚀 SINH ĐỀ V5.0 + HỘI ĐỒNG 3 AI", type="primary", use_container_width=True):
        if not selected_model: st.error("Chưa có mô hình AI.")
        else:
            try:
                bp={"total_questions":int(total_q),"type_distribution":{"mcq":int(mcq_n),"true_false":int(tf_n),"short_answer":int(short_n)},"level_distribution":{}}
                for item in levels.split(','):
                    if ':' in item:
                        k,v=item.split(':',1); bp["level_distribution"][k.strip()]=int(v.strip())
                validate_blueprint_counts(int(total_q), bp["type_distribution"], bp["level_distribution"])
                bp["topic_distribution"]=balanced_topic_distribution(topics, int(total_q))
                source_text=""
                if uploaded and uploaded.name.lower().endswith((".txt",".docx")):
                    source_text=read_source(uploaded)[0]
                model=genai.GenerativeModel(selected_model,generation_config={"response_mime_type":"application/json","temperature":0.2})
                with st.status("Đang chạy dây chuyền V5.0…", expanded=True) as status:
                    st.write("1/5 AI đang sinh đề theo ma trận…")
                    exam=parse_ai_json(model.generate_content(exam_generation_prompt(bp,source_text)).text)
                    st.write("2/5 Math Engine kiểm chứng…")
                    mr=audit_exam(exam)
                    st.write("3/5 Pedagogy Engine thẩm định…")
                    pr=audit_pedagogy(exam)
                    council=[]
                    for role,label in [("math","AI Toán học"),("pedagogy","AI Sư phạm"),("adversarial","AI Red-Team")]:
                        st.write(f"4/5 {label} phản biện…")
                        try: council.append(parse_ai_json(model.generate_content(reviewer_prompt(role,exam,mr,pr)).text))
                        except Exception as e: council.append({"role":role,"status":"REVIEW","score":0,"findings":[{"severity":"REVIEW","finding":f"Không đọc được phản biện: {e}"}]})
                    st.write("5/5 QA Gate + chứng nhận…")
                    rep=combined_v4_report(exam,mr,pr,council)
                    adaptive=analyze_exam(exam)
                    variants=build_variants(exam,int(variant_n),seed=1000)
                    manifest=build_manifest(variants,rep)
                    manifest["adaptive_analysis"]=adaptive
                    st.session_state["v4_exam"]=exam; st.session_state["v4_report"]=rep; st.session_state["v45_variants"]=variants; st.session_state["v45_manifest"]=manifest
                    status.update(label="Hoàn tất dây chuyền V5.0",state="complete",expanded=False)
            except Exception as e: st.error(f"V5.0 gặp lỗi: {e}")
    exam=st.session_state.get("v4_exam"); rep=st.session_state.get("v4_report")
    if exam and rep:
        c1,c2,c3=st.columns(3); c1.metric("QA Score",rep["summary"]["score"]); c2.metric("Trạng thái",rep["status"]); c3.metric("Chứng nhận",rep["certificate"]["certificate_status"])
        if rep["status"]=="PASS": st.success("🟢 ĐẠT KIỂM TRA TỰ ĐỘNG V5.0 — vẫn cần giáo viên duyệt cuối.")
        elif rep["status"]=="FAIL": st.error("🔴 REJECTED — còn lỗi FAIL, không nên phát hành.")
        else: st.warning("🟡 CONDITIONAL — cần giáo viên duyệt các điểm REVIEW.")
        for c in rep["council"]:
            with st.expander(f"{c.get('role','Reviewer')} — {c.get('status','REVIEW')} — {c.get('score',0)}"):
                st.json(c)
        with st.expander("🔐 Chứng nhận & dấu vết phiên bản"): st.json(rep["certificate"])
        st.download_button("📥 JSON báo cáo V5.0",json.dumps({"exam":exam,"report":rep},ensure_ascii=False,indent=2),"bao_cao_V5_0.json","application/json",use_container_width=True)
        st.download_button("📝 Xuất Word đề",build_exam_docx(exam,rep["certificate"]),"de_toan_V5_0.docx","application/vnd.openxmlformats-officedocument.wordprocessingml.document",use_container_width=True)
        st.download_button("📊 Xuất PowerPoint đề",build_exam_pptx(exam,theme_name,False),"de_toan_V5_0.pptx","application/vnd.openxmlformats-officedocument.presentationml.presentation",use_container_width=True)
        variants=st.session_state.get("v45_variants",[exam]); manifest=st.session_state.get("v45_manifest",{})
        st.markdown("### 🧠 Adaptive Intelligence")
        ad=manifest.get("adaptive_analysis",{})
        ac1,ac2,ac3=st.columns(3); ac1.metric("Độ khó heuristic TB",ad.get("summary",{}).get("avg_difficulty",0)); ac2.metric("Biên độ độ khó",ad.get("summary",{}).get("spread",0)); ac3.metric("Số mã đề",len(variants))
        st.caption("Điểm độ khó là heuristic hỗ trợ biên tập, không phải chỉ số tâm trắc học.")
        vc=variant_consistency(variants); st.write(f"**QA liên mã:** {vc['status']} — {len(vc.get('issues',[]))} cảnh báo")
        if vc.get("issues"): st.dataframe(vc["issues"],use_container_width=True)
        st.download_button("📦 Tải manifest + QA nhiều mã",json.dumps(manifest,ensure_ascii=False,indent=2),"manifest_v5_0.json","application/json",use_container_width=True)
    st.stop()

if mode in {"Thẩm định đề Toán Pro", "Thẩm định đề Toán 360°"}:
    st.subheader("2. Hệ thống thẩm định đề Toán 360°")
    st.info("V5.0: Math Engine + kiểm tra cấu trúc + ma trận + miền xác định + trùng/gần trùng + chất lượng sư phạm + QA Gate. Các tiêu chí heuristic chỉ đưa ra REVIEW, không tự kết luận thay giáo viên.")
    json_text = st.text_area("Dán trực tiếp JSON của đề", height=220, placeholder='{"questions": [...], "blueprint": {...}}')
    if st.button("🛡️ CHẠY THẨM ĐỊNH 360°", type="primary", use_container_width=True) and (uploaded or json_text.strip()):
        try:
            if json_text.strip(): exam=json.loads(json_text)
            else: exam=json.loads(uploaded.getvalue().decode("utf-8"))
            math_report=audit_exam(exam)
            ped_report=audit_pedagogy(exam)
            # Conservative combined gate: any deterministic FAIL keeps the exam FAIL.
            all_issues=math_report["exam_issues"]+ped_report["exam_issues"]
            combined_status="FAIL" if math_report["status"]=="FAIL" or ped_report["status"]=="FAIL" else "MANUAL_REVIEW" if math_report["status"]=="MANUAL_REVIEW" or ped_report["status"]=="MANUAL_REVIEW" else "PASS"
            combined_score=round((math_report["summary"]["score"]+ped_report["summary"]["score"])/2,1)
            st.session_state["exam_current"]=exam
            st.session_state["audit_report"]={"status":combined_status,"summary":{"total":len(exam.get("questions",[])),"pass":math_report["summary"]["pass"],"fail":max(math_report["summary"]["fail"],ped_report["summary"]["fail"]),"manual_review":max(math_report["summary"]["manual_review"],ped_report["summary"]["manual_review"]),"score":combined_score},"math":math_report,"pedagogy":ped_report,"exam_issues":all_issues}
        except Exception as e:
            st.error(f"Không đọc được đề JSON: {e}")
    report=st.session_state.get("audit_report")
    exam=st.session_state.get("exam_current")
    if report:
        sm=report["summary"]
        c1,c2,c3,c4,c5=st.columns(5)
        c1.metric("Tổng câu",sm["total"]); c2.metric("PASS",sm["pass"]); c3.metric("FAIL",sm["fail"]); c4.metric("REVIEW",sm["manual_review"]); c5.metric("360 Score",sm["score"])
        if report["status"]=="PASS": st.success("🟢 ĐỀ ĐẠT – qua cả kiểm tra Toán và kiểm tra sư phạm tự động.")
        elif report["status"]=="FAIL": st.error("🔴 ĐỀ KHÔNG ĐẠT – phải sửa lỗi FAIL trước khi phát hành.")
        else: st.warning("🟡 ĐỀ CẦN DUYỆT – có tiêu chí cần giáo viên thẩm định thủ công.")
        tab1,tab2,tab3=st.tabs(["🧮 Math Engine","🎓 Sư phạm 360°","📋 QA Gate"])
        with tab1:
            mr=report["math"]
            st.write(f"Trạng thái: **{mr['status']}** · QA Score: **{mr['summary']['score']}**")
            if mr["exam_issues"]: st.dataframe(mr["exam_issues"],use_container_width=True)
            table=[]
            for r in mr["questions"]: table.append({"Câu":r["number"],"Mức độ":r["level"],"Chủ đề":r["topic"],"Trạng thái":r["status"],"Số lỗi":len(r["issues"])})
            st.dataframe(table,use_container_width=True)
            with st.expander("🔎 Chi tiết Math Engine"):
                for r in mr["questions"]:
                    st.markdown(f"**Câu {r['number']} — {r['status']}**")
                    for x in r["issues"]: st.write(f"• {x['severity']} | {x['code']}: {x['message']} {x['evidence']}")
        with tab2:
            pr=report["pedagogy"]
            st.write(f"Trạng thái: **{pr['status']}** · Pedagogy Score: **{pr['summary']['score']}**")
            if pr["exam_issues"]: st.dataframe(pr["exam_issues"],use_container_width=True)
            ptable=[]
            for r in pr["questions"]: ptable.append({"Câu":r["number"],"Loại":r["type"],"Mức độ":r["level"],"Chủ đề":r["topic"],"Trạng thái":r["status"],"Số cảnh báo":len(r["issues"])})
            st.dataframe(ptable,use_container_width=True)
            with st.expander("🔎 Chi tiết kiểm tra sư phạm"):
                for r in pr["questions"]:
                    st.markdown(f"**Câu {r['number']} — {r['status']}**")
                    for x in r["issues"]: st.write(f"• {x['severity']} | {x['code']}: {x['message']} {x['evidence']}")
            st.caption("Lưu ý: kiểm tra mức độ nhận thức là heuristic hỗ trợ giáo viên; không phải phép đo tâm lý học hay thay thế thẩm định chuyên môn.")
        with tab3:
            st.markdown("**Nguyên tắc QA Gate:** FAIL = chặn phát hành; REVIEW = cần giáo viên duyệt; PASS = không phát hiện lỗi trong phạm vi engine.")
            if report["exam_issues"]: st.dataframe(report["exam_issues"],use_container_width=True)
            st.json({"status":report["status"],"summary":report["summary"]})
        if st.button("🔧 TỰ SỬA LỖI KỸ THUẬT + THẨM ĐỊNH 360° LẦN 2", use_container_width=True):
            fixed,changes=safe_autofix(exam)
            new_math=audit_exam(fixed); new_ped=audit_pedagogy(fixed)
            combined_status="FAIL" if new_math["status"]=="FAIL" or new_ped["status"]=="FAIL" else "MANUAL_REVIEW" if new_math["status"]=="MANUAL_REVIEW" or new_ped["status"]=="MANUAL_REVIEW" else "PASS"
            combined={"status":combined_status,"summary":{"total":len(fixed.get("questions",[])),"pass":new_math["summary"]["pass"],"fail":max(new_math["summary"]["fail"],new_ped["summary"]["fail"]),"manual_review":max(new_math["summary"]["manual_review"],new_ped["summary"]["manual_review"]),"score":round((new_math["summary"]["score"]+new_ped["summary"]["score"])/2,1)},"math":new_math,"pedagogy":new_ped,"exam_issues":new_math["exam_issues"]+new_ped["exam_issues"]}
            st.session_state["exam_current"]=fixed; st.session_state["audit_report"]=combined
            st.success(f"Đã sửa {len(changes)} lỗi kỹ thuật xác định chắc chắn.")
            if changes: st.code("\n".join(changes))
            st.rerun()
        payload={"version":"5.0.0","exam":exam,"audit":report}
        st.download_button("📥 Tải báo cáo thẩm định 360° JSON",json.dumps(payload,ensure_ascii=False,indent=2),"bao_cao_tham_dinh_360_v5_0.json","application/json",use_container_width=True)
    st.stop()

st.subheader("2. Lesson Studio V7 — Math Premium")
st.caption("Tạo cấu trúc → kiểm định → xem trước/chỉnh sửa → xuất PowerPoint.")
if st.button("🚀 TẠO CẤU TRÚC BÀI GIẢNG", type="primary", use_container_width=True, disabled=uploaded is None):
    try:
        config = LessonConfig(teacher, school, grade, book, lesson, int(periods), student_level,
                              int(slide_count), theme_name, include_answers, include_notes)
        source_text, source_bytes, source_type = read_source(uploaded)
        
        with st.status("Đang soạn bài giảng…", expanded=True) as status:
            st.write("Đang đọc và cấu trúc hóa tài liệu nguồn…")
            lesson_data = generate_lesson(selected_model, source_text, source_bytes, source_type, config)
            
            st.write("Đang kiểm định tiến trình, mật độ chữ và độ đa dạng bố cục…")
            lesson_report=audit_lesson(lesson_data,int(slide_count),source_text)
            st.session_state["lesson_data_v6"]=lesson_data
            st.session_state["lesson_config_v6"]=config
            st.session_state["lesson_source_v6"]=source_text
            st.session_state["lesson_report_v6"]=lesson_report
            status.update(label="Đã tạo cấu trúc — mời thầy duyệt trước khi xuất", state="complete", expanded=False)
                           
    except json.JSONDecodeError:
        st.error("AI trả về dữ liệu chưa đúng định dạng. Vui lòng bấm tạo lại.")
    except ValueError as exc:
        st.error(str(exc))
    except Exception as e:
        st.error(f"Lỗi: {str(e)}")

lesson_data=st.session_state.get("lesson_data_v6")
config=st.session_state.get("lesson_config_v6")
report=st.session_state.get("lesson_report_v6")
if lesson_data and config and report:
    st.markdown("### 3. Kiểm định và xem trước")
    a,b,c,d=st.columns(4)
    a.metric("Slide nội dung",report["summary"]["slides"])
    b.metric("Điểm QA",report["score"])
    c.metric("Trạng thái",report["status"])
    d.metric("Kiểu bố cục",len(report["summary"]["layouts"]))
    if report["status"]=="FAIL": st.error("Bài giảng thiếu thành phần bắt buộc; cần duyệt và sửa trước khi dùng.")
    elif report["status"]=="REVIEW": st.warning("Bài giảng đã tạo nhưng còn điểm cần giáo viên duyệt.")
    else: st.success("Không phát hiện lỗi trong phạm vi kiểm định tự động.")
    if report["issues"]: st.dataframe(report["issues"],use_container_width=True)
    preview=[{"STT":i,"Hoạt động":s["activity"],"Layout":s["layout"],"Tiêu đề":s["title"],"Nhiệm vụ":s.get("question","")[:90],"Sản phẩm":s.get("product","")[:70]} for i,s in enumerate(lesson_data["slides"],1)]
    st.dataframe(preview,use_container_width=True,height=360)
    with st.expander("✏️ Chỉnh sửa JSON bài giảng trước khi xuất"):
        edited=st.text_area("Nội dung cấu trúc",json.dumps(lesson_data,ensure_ascii=False,indent=2),height=440,key="lesson_json_editor")
        if st.button("✅ ÁP DỤNG CHỈNH SỬA VÀ KIỂM ĐỊNH LẠI",use_container_width=True):
            try:
                updated=validate_lesson(json.loads(edited))
                new_report=audit_lesson(updated,int(config.slide_count),st.session_state.get("lesson_source_v6",""))
                st.session_state["lesson_data_v6"]=updated; st.session_state["lesson_report_v6"]=new_report
                st.rerun()
            except Exception as exc: st.error(f"Không thể áp dụng: {exc}")
    st.markdown("### 4. Xuất PowerPoint")
    try:
        pptx_bytes=build_pptx(lesson_data,config)
        safe_name=re.sub(r"[^0-9A-Za-zÀ-ỹ_-]+","_",lesson_data.get("title") or "Bai_giang_Toan").strip("_")
        filename=f"{safe_name[:70]}_LessonStudioV7_MathPremium.pptx"
        st.download_button("📥 TẢI POWERPOINT MATH PREMIUM V7",pptx_bytes,filename,"application/vnd.openxmlformats-officedocument.presentationml.presentation",use_container_width=True)
        st.download_button("📋 TẢI BÁO CÁO KIỂM ĐỊNH",json.dumps(report,ensure_ascii=False,indent=2),f"{safe_name[:70]}_QA.json","application/json",use_container_width=True)
    except Exception as exc: st.error(f"Không thể dựng PowerPoint: {exc}")
